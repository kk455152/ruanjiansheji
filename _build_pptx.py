"""Generate the weekly acceptance PPTX for the smart-speaker mini-program project.

Layout follows the requirements in ppt制作要求.txt; visual style is inspired by
the dark navy + gold palette of pptx example.pdf.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ---------- Palette: deep navy + gold + ice blue (Midnight Executive) ----------
NAVY = RGBColor(0x0F, 0x1B, 0x3D)        # primary background
NAVY_DEEP = RGBColor(0x07, 0x10, 0x28)   # darker accent
GOLD = RGBColor(0xC9, 0xA0, 0x4E)        # accent
GOLD_LIGHT = RGBColor(0xE6, 0xC0, 0x7A)
ICE = RGBColor(0xCA, 0xDC, 0xFC)         # secondary
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF5, 0xF5, 0xF5)
MUTED = RGBColor(0x9A, 0xA8, 0xC2)
CARD = RGBColor(0x18, 0x27, 0x4D)

CN_HEADER_FONT = "Microsoft YaHei"
CN_BODY_FONT = "Microsoft YaHei"
EN_FONT = "Calibri"

IMG_DIR = r"c:\Users\32545\Desktop\软件设计\ruanjiansheji\_pptx_imgs"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height

blank_layout = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(blank_layout)


def fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def stroke(shape, rgb, width_pt=1):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def no_line(shape):
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, rgb, line=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    fill(s, rgb)
    if not line:
        no_line(s)
    return s


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=WHITE,
             font=CN_BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
        # also set east-asian font
        rPr = r._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', font)
    return tb


def slide_bg_navy(slide):
    add_rect(slide, 0, 0, SW, SH, NAVY)


def header_bar(slide, page_label, title_cn, title_en=None):
    # left vertical accent bar
    add_rect(slide, Inches(0.5), Inches(0.55), Inches(0.08), Inches(0.7), GOLD)
    # page label small
    add_text(slide, Inches(0.7), Inches(0.5), Inches(2.0), Inches(0.4),
             page_label, size=11, color=GOLD, bold=True, font=EN_FONT)
    # title big
    add_text(slide, Inches(0.7), Inches(0.78), Inches(11), Inches(0.7),
             title_cn, size=28, bold=True, color=WHITE, font=CN_HEADER_FONT)
    if title_en:
        add_text(slide, Inches(0.7), Inches(1.35), Inches(11), Inches(0.4),
                 title_en, size=11, color=MUTED, font=EN_FONT, italic=True)
    # bottom thin gold strip
    add_rect(slide, 0, SH - Inches(0.18), SW, Inches(0.18), NAVY_DEEP)
    add_rect(slide, 0, SH - Inches(0.04), SW, Inches(0.04), GOLD)


def page_number(slide, n, total):
    add_text(slide, SW - Inches(1.5), SH - Inches(0.55), Inches(1.2), Inches(0.35),
             f"{n:02d} / {total:02d}", size=10, color=MUTED, font=EN_FONT,
             align=PP_ALIGN.RIGHT)


def footer_brand(slide):
    add_text(slide, Inches(0.5), SH - Inches(0.55), Inches(6), Inches(0.35),
             "智能音箱音乐小程序  ·  本周验收汇报", size=10, color=MUTED,
             font=CN_BODY_FONT)


def add_image_clip(slide, path, x, y, w, h):
    """Add an image fitted into a box (padded inside, centered)."""
    if not os.path.exists(path):
        return None
    from PIL import Image
    im = Image.open(path)
    iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        # width-limited
        nw = w
        nh = int(w / img_ratio)
    else:
        nh = h
        nw = int(h * img_ratio)
    nx = x + (w - nw) // 2
    ny = y + (h - nh) // 2
    return slide.shapes.add_picture(path, nx, ny, nw, nh)


def img(name):
    return os.path.join(IMG_DIR, name)


def card(slide, x, y, w, h):
    s = add_rect(slide, x, y, w, h, CARD)
    return s


# ---------- Total slide count -----------------------------------------
TOTAL = 11

# =====================================================================
# SLIDE 1 — Title / cover
# =====================================================================
s1 = add_slide()
slide_bg_navy(s1)

# decorative diagonal gold band on the right
band = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.6), Inches(0),
                           Inches(0.06), SH)
fill(band, GOLD)

# small watermark text top-right
add_text(s1, Inches(9.5), Inches(0.6), Inches(3.5), Inches(0.4),
         "WEEKLY ACCEPTANCE REVIEW", size=10, color=GOLD, bold=True,
         font=EN_FONT, align=PP_ALIGN.RIGHT)
add_text(s1, Inches(9.5), Inches(0.95), Inches(3.5), Inches(0.4),
         "2026 / Week Report", size=10, color=MUTED,
         font=EN_FONT, align=PP_ALIGN.RIGHT)

# big english label
add_text(s1, Inches(0.9), Inches(1.6), Inches(8), Inches(0.7),
         "SMART  SPEAKER  ·  WEEKLY  REVIEW", size=14, color=GOLD, bold=True,
         font=EN_FONT)

# big Chinese title
add_text(s1, Inches(0.9), Inches(2.1), Inches(11), Inches(1.4),
         "智能音箱音乐小程序", size=58, bold=True, color=WHITE,
         font=CN_HEADER_FONT)
add_text(s1, Inches(0.9), Inches(3.3), Inches(11), Inches(0.9),
         "本周验收汇报", size=40, bold=True, color=GOLD_LIGHT,
         font=CN_HEADER_FONT)

# tagline
add_text(s1, Inches(0.9), Inches(4.4), Inches(10), Inches(0.5),
         "面向“优档”验收标准的设计 · 开发 · 测试 全流程汇报", size=15,
         color=ICE, italic=True, font=CN_BODY_FONT)

# accent line + team
add_rect(s1, Inches(0.9), Inches(5.3), Inches(0.6), Inches(0.04), GOLD)
add_text(s1, Inches(0.9), Inches(5.4), Inches(11), Inches(0.45),
         "汇报小组：A · B · C", size=14, color=WHITE, bold=True,
         font=CN_BODY_FONT)
add_text(s1, Inches(0.9), Inches(5.85), Inches(11), Inches(0.4),
         "课程：软件设计     周次：第 X 周     日期：2026.05.14",
         size=12, color=MUTED, font=CN_BODY_FONT)

# bottom-left badge
badge = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(6.6),
                            Inches(0.55), Inches(0.55))
fill(badge, GOLD)
add_text(s1, Inches(0.5), Inches(6.6), Inches(0.55), Inches(0.55),
         "01", size=14, bold=True, color=NAVY, font=EN_FONT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s1, Inches(1.2), Inches(6.65), Inches(8), Inches(0.45),
         "Cover  ·  项目封面", size=12, color=ICE, font=EN_FONT)

# =====================================================================
# SLIDE 2 — 本周验收标准与优档目标
# =====================================================================
s2 = add_slide()
slide_bg_navy(s2)
header_bar(s2, "02  ACCEPTANCE  CRITERIA", "本周验收标准与优档目标",
           "Weekly criteria & the 'Excellent' grade requirements")

# left: target box (large)
card(s2, Inches(0.5), Inches(1.85), Inches(4.6), Inches(5.0))
add_text(s2, Inches(0.8), Inches(2.05), Inches(4), Inches(0.4),
         "GOAL", size=11, bold=True, color=GOLD, font=EN_FONT)
add_text(s2, Inches(0.8), Inches(2.4), Inches(4), Inches(0.6),
         "本周目标：冲击「优」档", size=22, bold=True, color=WHITE,
         font=CN_HEADER_FONT)
add_rect(s2, Inches(0.8), Inches(3.05), Inches(0.5), Inches(0.04), GOLD)
add_text(s2, Inches(0.8), Inches(3.2), Inches(4), Inches(3.6),
         "PPT 不只展示「我们做了什么」，\n而要逐条对应「优」档的硬性要求。\n\n本汇报围绕 5 项验收点，\n用文档 + 截图 + 演示进行佐证，\n并在最后一页给出对照表。",
         size=14, color=ICE, font=CN_BODY_FONT, line_spacing=1.5)

# right: 5 criteria as numbered rows
criteria = [
    ("01", "有相关设计与开发文档", "Apifox 文档 · 数据库设计 · 接口字段表"),
    ("02", "小程序功能完整", "登录 · 设备 · 音乐/统计 三页面闭环"),
    ("03", "界面美观", "小程序运行截图 · 设计草图对照"),
    ("04", "实现与设计相符", "Apifox 字段 ↔ 前端字段对照表"),
    ("05", "MongoDB 与 MySQL 分类存储", "静态数据 / 动态数据 双库分流方案"),
]
y0 = Inches(1.85)
row_h = Inches(0.95)
for i, (num, title, desc) in enumerate(criteria):
    y = y0 + row_h * i + Emu(70000) * i
    # number circle
    o = s2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.4), y + Inches(0.15),
                            Inches(0.65), Inches(0.65))
    fill(o, GOLD)
    add_text(s2, Inches(5.4), y + Inches(0.15), Inches(0.65), Inches(0.65),
             num, size=14, bold=True, color=NAVY, font=EN_FONT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # title
    add_text(s2, Inches(6.2), y + Inches(0.05), Inches(6.7), Inches(0.5),
             title, size=18, bold=True, color=WHITE, font=CN_HEADER_FONT)
    # desc
    add_text(s2, Inches(6.2), y + Inches(0.5), Inches(6.7), Inches(0.4),
             desc, size=12, color=MUTED, font=CN_BODY_FONT)

footer_brand(s2)
page_number(s2, 2, TOTAL)

# =====================================================================
# SLIDE 3 — 小程序整体功能架构图
# =====================================================================
s3 = add_slide()
slide_bg_navy(s3)
header_bar(s3, "03  ARCHITECTURE", "小程序整体功能架构图",
           "Overall functional architecture of the mini-program")

# left description
add_text(s3, Inches(0.5), Inches(1.85), Inches(4.5), Inches(0.45),
         "三层闭环", size=18, bold=True, color=GOLD, font=CN_HEADER_FONT)
add_text(s3, Inches(0.5), Inches(2.3), Inches(4.5), Inches(3.5),
         "前端：微信小程序\n  · 登录 / 设备 / 音乐 / 统计\n\n"
         "中层：Flask + Nginx 后端\n  · RESTful API 网关\n  · 鉴权 · 转发 · 校验\n\n"
         "存储：MySQL + MongoDB\n  · 静态结构化数据\n  · 实时动态数据",
         size=13, color=ICE, font=CN_BODY_FONT, line_spacing=1.55)

# right: backend overview image
card(s3, Inches(5.2), Inches(1.85), Inches(7.6), Inches(5.0))
add_image_clip(s3, img("backend_overview.png"),
               Inches(5.35), Inches(2.0), Inches(7.3), Inches(4.7))
add_text(s3, Inches(5.2), Inches(6.5), Inches(7.6), Inches(0.35),
         "图 1  ·  后端整体框架与小程序的关系",
         size=10, color=MUTED, italic=True, font=CN_BODY_FONT,
         align=PP_ALIGN.CENTER)

footer_brand(s3)
page_number(s3, 3, TOTAL)

# =====================================================================
# SLIDE 4 — 小程序页面展示
# =====================================================================
s4 = add_slide()
slide_bg_navy(s4)
header_bar(s4, "04  MINI-PROGRAM PAGES", "小程序页面展示：登录 / 设备 / 音乐与统计",
           "Login · Device · Music & Stats — three-page closed loop")

screens = [
    ("mp_run1.png", "登录页", "微信授权 · Token 获取"),
    ("mp_run2.png", "设备页", "设备列表 · 绑定 · 设置"),
    ("mp_screen1.png", "音乐页", "搜索 · 播放 · 音量"),
    ("mp_screen2.png", "统计页", "周听歌数据 · 分享卡片"),
]
n = len(screens)
gap = Inches(0.25)
total_w = Inches(12.3)
each_w = Emu(int((total_w - gap * (n - 1)) / n))
x0 = Inches(0.5)
y0 = Inches(1.85)
phone_h = Inches(4.0)
for i, (fn, title, desc) in enumerate(screens):
    x = x0 + (each_w + gap) * i
    # phone frame card
    card(s4, x, y0, each_w, phone_h)
    add_image_clip(s4, img(fn), x + Inches(0.15), y0 + Inches(0.15),
                   each_w - Inches(0.3), phone_h - Inches(0.3))
    # title under card
    add_text(s4, x, y0 + phone_h + Inches(0.15), each_w, Inches(0.45),
             title, size=18, bold=True, color=WHITE, font=CN_HEADER_FONT,
             align=PP_ALIGN.CENTER)
    add_text(s4, x, y0 + phone_h + Inches(0.6), each_w, Inches(0.4),
             desc, size=11, color=MUTED, font=CN_BODY_FONT,
             align=PP_ALIGN.CENTER)

footer_brand(s4)
page_number(s4, 4, TOTAL)

# =====================================================================
# SLIDE 5 — Apifox 接口设计截图与接口列表
# =====================================================================
s5 = add_slide()
slide_bg_navy(s5)
header_bar(s5, "05  API DESIGN", "Apifox 接口设计与核心接口列表",
           "Apifox documentation & primary API endpoints")

# left: Apifox screenshot card
card(s5, Inches(0.5), Inches(1.85), Inches(6.4), Inches(5.0))
add_image_clip(s5, img("apifox_screen.png"),
               Inches(0.65), Inches(2.0), Inches(6.1), Inches(4.7))
add_text(s5, Inches(0.5), Inches(6.5), Inches(6.4), Inches(0.35),
         "图 2  ·  Apifox 接口工作台",
         size=10, color=MUTED, italic=True, font=CN_BODY_FONT,
         align=PP_ALIGN.CENTER)

# right: API list with mono-style cards
add_text(s5, Inches(7.2), Inches(1.85), Inches(5.8), Inches(0.4),
         "核心 API（节选）", size=16, bold=True, color=GOLD,
         font=CN_HEADER_FONT)
apis = [
    ("GET",  "/api/home/overview",       "获取首页信息"),
    ("POST", "/api/player/play-control", "播放控制"),
    ("POST", "/api/device/volume",       "音量调节"),
    ("GET",  "/api/friend/listening",    "听歌好友入口"),
    ("GET",  "/api/stats/weekly",        "听歌数据入口"),
    ("GET",  "/api/device/list",         "设备管理入口"),
    ("GET",  "/api/play-history",        "播放历史入口"),
    ("POST", "/api/music-service/bind",  "音乐平台绑定"),
]
y = Inches(2.35)
row_h = Inches(0.52)
for i, (m, path, label) in enumerate(apis):
    yy = y + row_h * i
    # method tag
    tag = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.2), yy + Inches(0.06),
                              Inches(0.7), Inches(0.32))
    fill(tag, GOLD if m == "GET" else ICE)
    no_line(tag)
    add_text(s5, Inches(7.2), yy + Inches(0.06), Inches(0.7), Inches(0.32),
             m, size=10, bold=True, color=NAVY, font=EN_FONT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # path
    add_text(s5, Inches(8.0), yy + Inches(0.02), Inches(3.4), Inches(0.4),
             path, size=11, color=WHITE, font="Consolas")
    # label
    add_text(s5, Inches(11.4), yy + Inches(0.02), Inches(1.7), Inches(0.4),
             label, size=10, color=MUTED, font=CN_BODY_FONT)

footer_brand(s5)
page_number(s5, 5, TOTAL)

# =====================================================================
# SLIDE 6 — MySQL 与 MongoDB 分类存储方案
# =====================================================================
s6 = add_slide()
slide_bg_navy(s6)
header_bar(s6, "06  STORAGE STRATEGY", "MySQL 与 MongoDB 分类存储方案",
           "Static structured data vs. real-time dynamic data")

# two-column comparison
col_w = Inches(6.0)
col_h = Inches(4.5)
col_y = Inches(1.95)

# MySQL column
card(s6, Inches(0.5), col_y, col_w, col_h)
add_rect(s6, Inches(0.5), col_y, col_w, Inches(0.5), GOLD)
add_text(s6, Inches(0.5), col_y, col_w, Inches(0.5),
         "MySQL  ·  关系型 · 静态结构化", size=14, bold=True, color=NAVY,
         font=CN_HEADER_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s6, Inches(0.7), col_y + Inches(0.7), Inches(5.6), Inches(4.2),
         "· user / auth_token  用户与登入 Token\n"
         "· device / user_device_binding  设备与绑定关系\n"
         "· media_mapping  歌曲元数据映射\n"
         "· play_history  播放历史\n"
         "· friendship  好友关系\n"
         "· listen_room / listen_room_member  一起听\n"
         "· share_record  分享行为\n"
         "· music_service_binding  第三方平台账号\n"
         "· operation_log / Daily_Stats  日志与统计\n"
         "· device_settings / battery_notice_setting",
         size=12, color=ICE, font=CN_BODY_FONT, line_spacing=1.5)

# MongoDB column
card(s6, Inches(6.85), col_y, col_w, col_h)
add_rect(s6, Inches(6.85), col_y, col_w, Inches(0.5), ICE)
add_text(s6, Inches(6.85), col_y, col_w, Inches(0.5),
         "MongoDB  ·  文档型 · 实时动态", size=14, bold=True, color=NAVY,
         font=CN_HEADER_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s6, Inches(7.05), col_y + Inches(0.7), Inches(5.6), Inches(4.2),
         "· device_realtime_status  设备实时状态\n"
         "    电量 / 在线 / 当前播放\n\n"
         "· play_event_stream  播放事件流\n"
         "    高频写入 · 灵活字段\n\n"
         "· voice_command_log  语音指令日志\n\n"
         "· user_session_cache  会话级缓存\n\n"
         "· stats_snapshot  统计中间结果",
         size=12, color=ICE, font=CN_BODY_FONT, line_spacing=1.5)

# bottom note
add_text(s6, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.35),
         "分流原则：结构稳定 · 关系强 → MySQL ；写入频繁 · 字段灵活 → MongoDB",
         size=11, italic=True, color=GOLD, font=CN_BODY_FONT,
         align=PP_ALIGN.CENTER)

footer_brand(s6)
page_number(s6, 6, TOTAL)

# =====================================================================
# SLIDE 7 — 后端接口与双数据库连接说明
# =====================================================================
s7 = add_slide()
slide_bg_navy(s7)
header_bar(s7, "07  BACKEND  ·  DATA FLOW", "后端接口与双数据库连接说明",
           "Flask backend orchestrating MySQL + MongoDB")

# left flow diagram screenshot
card(s7, Inches(0.5), Inches(1.85), Inches(7.5), Inches(5.0))
add_image_clip(s7, img("data_flow_arch.png"),
               Inches(0.65), Inches(2.0), Inches(7.2), Inches(4.7))
add_text(s7, Inches(0.5), Inches(6.5), Inches(7.5), Inches(0.35),
         "图 3  ·  核心数据流动与安全校验架构",
         size=10, color=MUTED, italic=True, font=CN_BODY_FONT,
         align=PP_ALIGN.CENTER)

# right: 4 step explainer
add_text(s7, Inches(8.3), Inches(1.85), Inches(4.7), Inches(0.4),
         "数据请求生命周期", size=16, bold=True, color=GOLD,
         font=CN_HEADER_FONT)
steps = [
    ("①", "小程序请求", "wx.request 发往 Nginx → Flask"),
    ("②", "鉴权 & 校验", "Token 校验 · 参数校验 · 权限分级"),
    ("③", "双库分流",   "MySQL 取静态 / MongoDB 取动态"),
    ("④", "逻辑拼装",   "组合静态 + 动态信息后返回前端"),
]
sy = Inches(2.4)
for i, (num, t, d) in enumerate(steps):
    y = sy + Inches(1.05) * i
    o = s7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.3), y,
                            Inches(0.55), Inches(0.55))
    fill(o, GOLD)
    add_text(s7, Inches(8.3), y, Inches(0.55), Inches(0.55), num,
             size=14, bold=True, color=NAVY, font=EN_FONT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s7, Inches(9.0), y - Inches(0.02), Inches(4.0), Inches(0.4),
             t, size=14, bold=True, color=WHITE, font=CN_HEADER_FONT)
    add_text(s7, Inches(9.0), y + Inches(0.42), Inches(4.0), Inches(0.5),
             d, size=11, color=MUTED, font=CN_BODY_FONT)

footer_brand(s7)
page_number(s7, 7, TOTAL)

# =====================================================================
# SLIDE 8 — 功能演示流程
# =====================================================================
s8 = add_slide()
slide_bg_navy(s8)
header_bar(s8, "08  END-TO-END FLOW", "功能演示流程：小程序 → 接口 → 数据库",
           "User action → API → DB → response, illustrated")

# Three big stage cards in a row, with arrows between them
stages = [
    ("用户操作",  "Mini-Program",
     "用户点击播放、调整音量、\n查看好友、查看周统计"),
    ("接口请求",  "Flask API",
     "携带 Token 调用 RESTful API\nApifox 已定义全部字段"),
    ("数据落库",  "MySQL + MongoDB",
     "静态写 MySQL · 动态写 MongoDB\n返回拼装结果给小程序"),
]
sx = Inches(0.5)
sy = Inches(2.3)
sw = Inches(3.9)
sh = Inches(3.1)
gap = Inches(0.3)
for i, (title_cn, title_en, body) in enumerate(stages):
    x = sx + (sw + gap) * i
    card(s8, x, sy, sw, sh)
    add_rect(s8, x, sy, sw, Inches(0.08), GOLD)
    add_text(s8, x + Inches(0.3), sy + Inches(0.3), sw, Inches(0.35),
             title_en, size=11, bold=True, color=GOLD, font=EN_FONT)
    add_text(s8, x + Inches(0.3), sy + Inches(0.65), sw - Inches(0.3),
             Inches(0.6), title_cn, size=22, bold=True, color=WHITE,
             font=CN_HEADER_FONT)
    add_text(s8, x + Inches(0.3), sy + Inches(1.45), sw - Inches(0.6),
             sh - Inches(1.6), body, size=13, color=ICE,
             font=CN_BODY_FONT, line_spacing=1.5)
    if i < 2:
        ax = x + sw + Emu(20000)
        ay = sy + Inches(1.2)
        arrow = s8.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax,
                                    ay, gap - Emu(40000), Inches(0.6))
        fill(arrow, GOLD)
        no_line(arrow)

# bottom row: live demo evidence images (2 small)
add_text(s8, Inches(0.5), Inches(5.7), Inches(12), Inches(0.35),
         "现场演示证据  Live demo evidence", size=12, bold=True,
         color=GOLD, font=CN_HEADER_FONT)
add_image_clip(s8, img("api_list_long.png"),
               Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.85))

footer_brand(s8)
page_number(s8, 8, TOTAL)

# =====================================================================
# SLIDE 9 — 测试截图
# =====================================================================
s9 = add_slide()
slide_bg_navy(s9)
header_bar(s9, "09  TESTS & EVIDENCE",
           "测试截图：Apifox · MySQL · MongoDB · 小程序运行",
           "Test artefacts proving end-to-end correctness")

# 2x2 grid of screenshots
boxes = [
    ("Apifox 接口测试",  "apifox_screen.png"),
    ("MySQL 表数据",     "mysql_table1.png"),
    ("MongoDB 文档",     "mongodb_doc1.png"),
    ("小程序运行",       "mp_run1.png"),
]
gx = Inches(0.5)
gy = Inches(1.95)
gw = Inches(6.1)
gh = Inches(2.45)
gap = Inches(0.25)
for i, (title, fn) in enumerate(boxes):
    col = i % 2
    row = i // 2
    x = gx + (gw + gap) * col
    y = gy + (gh + gap) * row
    card(s9, x, y, gw, gh)
    # caption strip
    add_rect(s9, x, y, gw, Inches(0.4), CARD)
    add_rect(s9, x, y, Inches(0.08), gh, GOLD)
    add_text(s9, x + Inches(0.25), y + Inches(0.05), gw - Inches(0.3),
             Inches(0.35), title, size=13, bold=True, color=WHITE,
             font=CN_HEADER_FONT)
    add_image_clip(s9, img(fn), x + Inches(0.15), y + Inches(0.5),
                   gw - Inches(0.3), gh - Inches(0.65))

footer_brand(s9)
page_number(s9, 9, TOTAL)

# =====================================================================
# SLIDE 10 — 验收标准逐条对照表 (the most critical slide)
# =====================================================================
s10 = add_slide()
slide_bg_navy(s10)
header_bar(s10, "10  ACCEPTANCE CHECKLIST",
           "验收标准逐条对照表",
           "Each criterion mapped to the evidence we produced")

# table
rows = [
    ("优档要求", "我们的证明材料", "对应页"),
    ("有相关设计和开发文档",
     "Apifox 文档 · 数据库设计表 · 接口字段表 · 设计与开发文档",
     "P3 · P5 · P6"),
    ("小程序功能完整",
     "登录 / 设备 / 音乐 / 统计 三页面闭环演示",
     "P4 · P8"),
    ("界面美观",
     "小程序运行截图 + 设计草图对照",
     "P4 · P9"),
    ("实现与设计相符",
     "Apifox 字段 ↔ 前端字段对照 · 端到端调用链",
     "P5 · P7 · P8"),
    ("MongoDB 与 MySQL 分类存储",
     "MySQL 表截图 + MongoDB 文档截图 + 分类说明",
     "P6 · P9"),
]

tx = Inches(0.5)
ty = Inches(2.0)
tw = Inches(12.3)
header_h = Inches(0.55)
row_h = Inches(0.78)

# header
add_rect(s10, tx, ty, tw, header_h, GOLD)
col_x = [tx, tx + Inches(3.6), tx + Inches(9.8)]
col_w = [Inches(3.6), Inches(6.2), Inches(2.5)]
for i, (cx, cw) in enumerate(zip(col_x, col_w)):
    add_text(s10, cx + Inches(0.25), ty, cw, header_h, rows[0][i],
             size=14, bold=True, color=NAVY, font=CN_HEADER_FONT,
             anchor=MSO_ANCHOR.MIDDLE)

# body rows
for ri in range(1, 6):
    y = ty + header_h + row_h * (ri - 1)
    bg = CARD if ri % 2 == 1 else NAVY_DEEP
    add_rect(s10, tx, y, tw, row_h, bg)
    # left accent
    add_rect(s10, tx, y, Inches(0.08), row_h, GOLD)
    for i, (cx, cw) in enumerate(zip(col_x, col_w)):
        text = rows[ri][i]
        is_first = (i == 0)
        add_text(s10, cx + Inches(0.25), y, cw, row_h, text,
                 size=13 if is_first else 12,
                 bold=is_first, color=WHITE if is_first else ICE,
                 font=CN_HEADER_FONT if is_first else CN_BODY_FONT,
                 anchor=MSO_ANCHOR.MIDDLE)

footer_brand(s10)
page_number(s10, 10, TOTAL)

# =====================================================================
# SLIDE 11 — 总结与后续优化 + Thank You
# =====================================================================
s11 = add_slide()
slide_bg_navy(s11)
header_bar(s11, "11  SUMMARY",
           "总结与后续优化",
           "Recap & next-week roadmap")

# left summary card
card(s11, Inches(0.5), Inches(1.95), Inches(6.0), Inches(4.5))
add_rect(s11, Inches(0.5), Inches(1.95), Inches(0.08), Inches(4.5), GOLD)
add_text(s11, Inches(0.8), Inches(2.15), Inches(5.5), Inches(0.5),
         "本周成果", size=20, bold=True, color=GOLD, font=CN_HEADER_FONT)
add_text(s11, Inches(0.8), Inches(2.7), Inches(5.5), Inches(3.7),
         "✓ 完成 5 项「优」档验收要求的全部佐证\n"
         "✓ 文档闭环：草图 → 接口 → 数据库 → 字段表\n"
         "✓ 小程序三页面联调可正常运行\n"
         "✓ MySQL + MongoDB 分流方案已落地\n"
         "✓ Apifox 文档与前端字段一一对应",
         size=14, color=ICE, font=CN_BODY_FONT, line_spacing=1.7)

# right next-steps card
card(s11, Inches(6.85), Inches(1.95), Inches(6.0), Inches(4.5))
add_rect(s11, Inches(6.85), Inches(1.95), Inches(0.08), Inches(4.5), ICE)
add_text(s11, Inches(7.15), Inches(2.15), Inches(5.5), Inches(0.5),
         "后续优化方向", size=20, bold=True, color=ICE, font=CN_HEADER_FONT)
add_text(s11, Inches(7.15), Inches(2.7), Inches(5.5), Inches(3.7),
         "• 接口性能压测（JMeter）\n"
         "• 引入 RabbitMQ 削峰填谷\n"
         "• MongoDB 索引与 TTL 策略优化\n"
         "• 完善权限分级与审计日志\n"
         "• 小程序 UI 细节打磨与可访问性",
         size=14, color=ICE, font=CN_BODY_FONT, line_spacing=1.7)

# bottom thank-you
add_text(s11, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.35),
         "THANK YOU  ·  感谢聆听  ·  欢迎指导",
         size=14, bold=True, color=GOLD, font=EN_FONT,
         align=PP_ALIGN.CENTER, italic=True)

footer_brand(s11)
page_number(s11, 11, TOTAL)

# =====================================================================
# Save
# =====================================================================
out_path = r"C:\Users\32545\Desktop\本周验收汇报.pptx"
prs.save(out_path)
print("Saved:", out_path)




